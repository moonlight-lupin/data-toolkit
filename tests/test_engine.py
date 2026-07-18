"""Regression suite for the data-toolkit's finance-grade behaviours.

Locks in exactly the guarantees that make the engine safe for finance work — the things that
would silently corrupt a reconciliation or a clean table if they regressed:

  - amounts are exact Decimal (no binary-float drift at the tolerance edge)
  - 100 USD never reconciles against 100 SGD (currencies are compared)
  - strict-currency mode refuses to match an unknown currency
  - a bare "$" is ambiguous, not assumed USD
  - amount/date pairs outside the window are flagged ambiguous, not matched
  - multi-tab workbooks require an explicit sheet (no silent 'active'-sheet guess)
  - form extraction handles next-line and dotted-leader layouts
  - a currency column can be split out via code_target
  - categorical value-map clustering proposes the right canonical
  - PDF tables: best-engine scoring (+ a guarded end-to-end smoke test)

Runs two ways:
    python tests/test_engine.py        # standalone, no pytest needed
    pytest tests/                       # if pytest is installed
"""

from __future__ import annotations

import sys
import tempfile
import json
import subprocess
from decimal import Decimal
from pathlib import Path

# --- make the shared engine + skill scripts importable, standalone or under pytest ---------
_ROOT = Path(__file__).resolve().parents[1]
for _p in (_ROOT / "scripts",
           _ROOT / "skills" / "data-reconcile" / "scripts",
           _ROOT / "skills" / "data-visualise" / "scripts",
           _ROOT / "skills" / "data-analyse" / "scripts"):
    if str(_p) not in sys.path:
        sys.path.insert(0, str(_p))

import dataclean        # noqa: E402
import extract          # noqa: E402
import ingest           # noqa: E402
import reconcile        # noqa: E402
import viz              # noqa: E402
import workbook         # noqa: E402


# --------------------------------------------------------------------------- #
# 1. Decimal tolerance edge
# --------------------------------------------------------------------------- #
def test_decimal_amounts_are_exact():
    n1, _ = dataclean.parse_number("0.1")
    n2, _ = dataclean.parse_number("0.2")
    assert isinstance(n1, Decimal)
    assert n1 + n2 == Decimal("0.3")                       # float would give 0.30000000000000004
    assert reconcile.to_amount("0.1") + reconcile.to_amount("0.2") == reconcile.to_amount("0.3")


def test_reconcile_tolerance_edge_is_exact():
    # a 0.10 spread with default tol 0.01 must be a value difference, not a fuzzy match,
    # and the reported diff must be the exact Decimal, not float dust.
    A = [{"k": "1", "amount": "1000.10"}]
    B = [{"k": "1", "amount": "1000.00"}]
    res = reconcile.match(A, B, key="k", amount="amount")
    assert not res["matched"] and len(res["value_diffs"]) == 1
    assert res["value_diffs"][0]["diff"] == Decimal("0.10")


# --------------------------------------------------------------------------- #
# 2. 100 USD vs 100 SGD must NOT reconcile (currency compared)
# --------------------------------------------------------------------------- #
def test_usd_does_not_match_sgd_on_key():
    A = [{"ref": "R1", "amount": "100.00", "ccy": "USD"}]
    B = [{"ref": "R1", "amount": "100.00", "ccy": "SGD"}]
    res = reconcile.match(A, B, key="ref", amount="amount", currency="ccy")
    assert not res["matched"] and len(res["currency_diffs"]) == 1
    assert any(e["category"] == "currency_mismatch" for e in reconcile.triage(res))


def test_usd_does_not_match_sgd_via_symbol_in_amount_date():
    # no currency column — code is detected from the amount cell's symbol
    res = reconcile.match([{"amount": "US$ 100"}], [{"amount": "S$ 100"}],
                          amount="amount", mode="amount_date")
    assert not res["matched"] and not res["ambiguous"]


# --------------------------------------------------------------------------- #
# 3. bare "$" is ambiguous (never assumed USD)
# --------------------------------------------------------------------------- #
def test_bare_dollar_is_ambiguous():
    (amt, code), note = dataclean.parse_currency("$1,000")
    assert amt == Decimal("1000") and code is None and "ambiguous" in note
    assert reconcile.to_currency("$100") is None
    # an expected currency resolves it
    val, note2, _ = dataclean._convert("$1,000", {"type": "currency", "currency": "SGD"})
    assert val == Decimal("1000") and note2 == ""


def test_disambiguated_dollars_resolve():
    assert dataclean.parse_currency("US$ 50")[0][1] == "USD"
    assert dataclean.parse_currency("S$ 2.5m")[0] == (Decimal("2500000.0"), "SGD")
    assert reconcile.to_currency("A$ 90") == "AUD"


# --------------------------------------------------------------------------- #
# 4. amount/date outside the window -> ambiguous, inside -> matched
# --------------------------------------------------------------------------- #
def test_amount_date_window_constraint():
    A = [{"amount": "500.00", "date": "01 Jun 2026"}]
    B = [{"amount": "500.00", "date": "25 Jun 2026"}]      # 24 days apart
    out = reconcile.match(A, B, amount="amount", date="date",
                          mode="amount_date", date_window_days=5)
    assert not out["matched"] and len(out["ambiguous"]) == 1
    assert any(e["category"] == "ambiguous_match" for e in reconcile.triage(out))
    inside = reconcile.match(A, B, amount="amount", date="date",
                             mode="amount_date", date_window_days=30)
    assert len(inside["matched"]) == 1 and not inside["ambiguous"]


# --------------------------------------------------------------------------- #
# 5. strict-currency mode refuses unknown-currency matches
# --------------------------------------------------------------------------- #
def test_strict_currency_blocks_unknown():
    A = [{"ref": "R1", "amount": "100.00"}]                # currency unknown
    B = [{"ref": "R1", "amount": "100.00"}]
    assert len(reconcile.match(A, B, key="ref", amount="amount")["matched"]) == 1   # permissive
    strict = reconcile.match(A, B, key="ref", amount="amount", strict_currency=True)
    assert not strict["matched"] and len(strict["currency_unknown"]) == 1
    assert any(e["category"] == "currency_unknown" for e in reconcile.triage(strict))


# --------------------------------------------------------------------------- #
# 6. multi-sheet workbook requires explicit selection
# --------------------------------------------------------------------------- #
def _write_xlsx(path, sheets):
    """sheets = [(name, rows_or_None)]; None leaves the sheet empty."""
    from openpyxl import Workbook
    wb = Workbook()
    wb.remove(wb.active)
    for name, rows in sheets:
        ws = wb.create_sheet(name)
        for r in (rows or []):
            ws.append(r)
    wb.save(path)


def test_multisheet_requires_selection():
    d = tempfile.mkdtemp()
    multi = Path(d) / "multi.xlsx"
    _write_xlsx(multi, [("Cover", None),
                        ("Data", [["a", "b"], [1, 2]]),
                        ("Notes", [["x"], ["hi"]])])
    try:
        ingest.read_any(str(multi))
        raise AssertionError("expected SheetSelectionRequired for a multi-data-sheet workbook")
    except ingest.SheetSelectionRequired as e:
        assert "Data" in str(e) and "Notes" in str(e)
    rows, note = ingest.read_any(str(multi), sheet="Data")
    assert rows[0] == ["a", "b"]

    single = Path(d) / "single.xlsx"
    _write_xlsx(single, [("Cover", None), ("TheData", [["k"], [9]])])
    rows, note = ingest.read_any(str(single))
    assert rows[0][0] == "k" and "auto-selected" in note      # picks the only data sheet


def test_viz_rows_from_xlsx_multisheet_safe():
    d = tempfile.mkdtemp()
    multi = Path(d) / "m.xlsx"
    _write_xlsx(multi, [("Cover", None),
                        ("Data", [["k", "v"], [1, 2]]),
                        ("More", [["x"], [9]])])
    try:
        viz.rows_from_xlsx(str(multi))
        raise AssertionError("expected a raise for a multi-data-sheet workbook")
    except Exception as e:                                   # SheetSelectionRequired or ValueError
        assert "sheet" in str(e).lower() or "Data" in str(e)
    assert viz.rows_from_xlsx(str(multi), sheet="Data") == [{"k": 1, "v": 2}]


def test_workbook_excel_charts_and_analysis_handoff():
    path = workbook.write_charts_xlsx(
        Path(tempfile.mkdtemp()) / "charts.xlsx",
        [
            {"chart_type": "column", "title": "By day",
             "categories": ["Mon", "Tue"], "series": [{"name": "Done", "values": [4, 7]}]},
            {"chart_type": "line", "title": "Trend",
             "categories": ["W1", "W2", "W3"],
             "series": [{"name": "Open", "values": [10, 14, 9]},
                        {"name": "Closed", "values": [8, 11, 13]}]},
            {"chart_type": "waterfall", "title": "Bridge",
             "steps": [
                 {"label": "Open", "value": 12, "kind": "start"},
                 {"label": "Done", "value": -5, "kind": "delta"},
                 {"label": "Close", "value": 7, "kind": "total"},
             ]},
        ],
        workbook_title="Self-test",
    )
    assert path.is_file() and path.stat().st_size > 1000
    from openpyxl import load_workbook
    wb = load_workbook(path)
    assert "By day" in wb.sheetnames and len(wb["By day"]._charts) == 1
    assert len(wb["Bridge"]._charts) == 1

    analysis = {
        "results": [
            {"op": "breakdown", "name": "Cust", "result": {
                "by": "Customer",
                "groups": [{"key": "A", "total": "100", "count": 1},
                           {"key": "B", "total": "40", "count": 1}],
            }},
            {"op": "period_series", "name": "Monthly", "result": {
                "grain": "month",
                "periods": [
                    {"period": "2026-01", "total": "100", "delta": None},
                    {"period": "2026-02", "total": "130", "delta": "30"},
                ],
            }},
            {"op": "concentration", "name": "skip", "result": {"hhi": 2000}},
        ]
    }
    specs = workbook.suggest_charts_from_analysis(analysis)
    types = [s["chart_type"] for s in specs]
    assert "column" in types and "pie" in types and "line" in types and "waterfall" in types
    assert all(s["chart_type"] != "concentration" for s in specs)
    out = workbook.charts_from_analysis(analysis, Path(tempfile.mkdtemp()) / "from-analysis.xlsx")
    assert out.is_file()


def test_officecli_render_is_optional_and_never_locks_the_workbook():
    """The OfficeCLI renderer is opt-in. Absent → degrade silently; present → PNGs *and* the
    workbook must be released (a resident holds an OS file lock until it is closed)."""
    import os

    import officecli_render as ocr

    charts = [{"chart_type": "column", "title": "T",
               "categories": ["A", "B"], "series": {"V": [1, 2]}}]

    # --- always: with the binary unavailable nothing raises and nothing is produced ---
    real_which = ocr.shutil.which
    try:
        ocr.shutil.which = lambda _name: None
        td = Path(tempfile.mkdtemp())
        src = workbook.write_charts_xlsx(td / "absent.xlsx", charts)
        assert ocr.available() is False
        assert ocr.chart_paths(src) == []
        assert ocr.render_chart_pngs(src, td / "out") == []
        assert ocr.render_sheet_png(src, td / "s.png") is None
        assert ocr.status() == {"available": False, "path": None, "version": None}
        assert src.is_file()                      # the workbook itself is untouched
    finally:
        ocr.shutil.which = real_which

    # --- only when actually installed: render, then prove the file is not left locked ---
    if not ocr.available():
        return
    td = Path(tempfile.mkdtemp())
    src = workbook.write_charts_xlsx(td / "live.xlsx", charts)
    pngs = ocr.render_chart_pngs(src, td / "out")
    assert pngs and all(p.is_file() and p.stat().st_size > 0 for p in pngs)
    probe = str(src) + ".probe"
    os.replace(src, probe)                        # raises PermissionError if still locked
    os.replace(probe, src)


def test_officecli_render_degrades_gracefully_on_every_failure_mode():
    """OfficeCLI is *optional*, so the contract is absolute: once the .xlsx is written,
    nothing the renderer does can fail the run. Every hostile mode must return a value —
    a non-zero exit, unparseable output, an unwritable destination, or an outright
    exception from the subprocess layer."""
    import officecli_render as ocr

    td = Path(tempfile.mkdtemp())
    src = td / "x.xlsx"
    src.write_bytes(b"not really a workbook")   # never opened: _run is stubbed throughout
    real_which, real_run = ocr.shutil.which, ocr._run
    try:
        ocr.shutil.which = lambda _name: "/fake/officecli"      # pretend it is installed
        assert ocr.available() is True

        ocr._run = lambda _a, timeout=60: (1, "", "boom")       # non-zero exit
        assert ocr.render_chart_pngs(src, td / "o1") == []

        ocr._run = lambda _a, timeout=60: (0, "<<<not json>>>", "")   # garbage stdout
        assert ocr.chart_paths(src) == []
        assert ocr.render_chart_pngs(src, td / "o2") == []

        # a *file* sits where out_dir must be created -> mkdir raises, we must not
        blocker = td / "o3"
        blocker.write_text("in the way", encoding="utf-8")
        ocr._run = lambda _a, timeout=60: (0, '{"data":{"results":[{"path":"/S/chart[1]"}]}}', "")
        assert ocr.render_chart_pngs(src, blocker) == []

        def explode(_a, timeout=60):                            # unexpected subprocess error
            raise RuntimeError("subprocess exploded")
        ocr._run = explode
        assert ocr.render_chart_pngs(src, td / "o4") == []
        assert ocr.render_sheet_png(src, td / "o5" / "s.png") is None
        assert ocr.chart_paths(src) == []
        assert ocr.version() is None
        assert ocr.release(src) is False          # runs inside a finally: must never raise
        assert ocr.status()["available"] is True  # probing still answers
    finally:
        ocr.shutil.which, ocr._run = real_which, real_run


def test_render_png_never_costs_the_run_its_workbook():
    """End-to-end: a plan asking for render_png must still yield the .xlsx when the optional
    renderer is missing, un-importable, or blows up — degraded to a warning, never an error."""
    import builtins

    import agent_runtime

    td = Path(tempfile.mkdtemp())
    src = td / "d.csv"
    src.write_text("region,amount\nNorth,10\nSouth,20\n", encoding="utf-8")
    plan = {
        "version": 1, "skill": "data-visualise", "input": str(src),
        "output": str(td / "c.xlsx"), "format": "xlsx",
        "dashboard": {"title": "By region", "render_png": True, "blocks": [
            {"type": "chart", "chart_type": "bar", "title": "By region",
             "categories": ["North", "South"],
             "series": [{"name": "Amount", "values": [10, 20]}]}]},
    }

    def assert_survives(label):
        env = agent_runtime.run_plan(json.loads(json.dumps(plan)), base_dir=td)
        assert env["status"] == "success_with_warnings", f"{label}: {env.get('errors')}"
        assert [a["kind"] for a in env["artifacts"]] == ["xlsx_charts"], label
        assert any("render_png" in w for w in env["warnings"]), label

    real_import = builtins.__import__

    def blocked(name, *a, **k):
        if name == "officecli_render":
            raise ImportError("simulated: module missing")
        return real_import(name, *a, **k)

    builtins.__import__ = blocked
    try:
        assert_survives("import fails")
    finally:
        builtins.__import__ = real_import

    import officecli_render as ocr
    real_avail, real_render = ocr.available, ocr.render_chart_pngs
    try:
        ocr.available = lambda: True

        def boom(*_a, **_k):
            raise RuntimeError("renderer exploded")
        ocr.render_chart_pngs = boom
        assert_survives("renderer raises")

        ocr.available = lambda: False
        assert_survives("binary absent")
    finally:
        ocr.available, ocr.render_chart_pngs = real_avail, real_render


def test_workbook_charts_follow_the_visualise_theme():
    """A white-label theme must colour the Excel workbook as it colours the HTML dashboard —
    one palette drives both artefacts (regression: the xlsx path used to ignore `theme`)."""
    import re
    import zipfile

    charts = [{"chart_type": "column", "title": "By region",
               "categories": ["N", "S"], "series": {"Amount": [10, 20]}}]
    out = Path(tempfile.mkdtemp())

    def series_fills(path):
        with zipfile.ZipFile(path) as z:
            xml = "".join(z.read(n).decode("utf8", "ignore")
                          for n in z.namelist() if "charts/chart" in n)
        return set(re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"', xml))

    default_fills = series_fills(workbook.write_charts_xlsx(out / "default.xlsx", charts))
    acme_fills = series_fills(workbook.write_charts_xlsx(
        out / "acme.xlsx", charts, theme={"colours": {"burgundy": "#0B3D91"}}))

    assert "163F3A" in default_fills, default_fills          # neutral default palette
    assert "0B3D91" in acme_fills, acme_fills                # the theme reached the chart
    assert "0B3D91" not in default_fills
    try:
        workbook.write_charts_xlsx(Path(tempfile.mkdtemp()) / "x.xlsx", [])
        raise AssertionError("empty charts should fail")
    except ValueError:
        pass


def _tooltip_counts(html):
    """Frequencies read back out of a histogram's per-bar tooltips."""
    import re
    return [int(m) for m in re.findall(r": (\d+)</title>", html)]


def test_viz_histogram_bins_and_forced_zero_axis():
    """Binning is where a histogram silently lies. Edges are half-open [lo,hi) except
    the last (which must include its upper bound, or the maximum vanishes), and the
    y-axis is forced to 0 — a floated frequency axis misstates every bar."""
    import re
    assert _tooltip_counts(viz.histogram([5, 35, 65, 120],
                                         bins=[0, 30, 60, 90, 365])) == [1, 1, 1, 1]
    assert _tooltip_counts(viz.histogram([5], bins=[0, 5, 10])) == [0, 1]   # half-open
    assert _tooltip_counts(viz.histogram([0, 10], bins=[0, 5, 10])) == [1, 1]  # last incl.
    assert sum(_tooltip_counts(viz.histogram([1, 1, 1, 2, 2, 3], bins=3))) == 6

    # a far-from-zero cluster must still show a 0 baseline
    hi = viz.histogram([100, 101, 102], bins=3)
    assert "0" in re.findall(r'class="x">([\d,.-]+)</text>', hi)

    # non-numeric vs out-of-range are DIFFERENT exclusions and are reported as such
    junk = viz.histogram([1, "abc", 3])
    assert "non-numeric" in junk
    out = viz.histogram([-5, 1, 999], bins=[0, 10])
    assert "outside the bin range" in out and "non-numeric" not in out

    assert "No numeric values" in viz.histogram(["a", ""])
    assert "two bin edges" in viz.histogram([1, 2], bins=[5])
    assert viz.histogram([7, 7, 7]).startswith('<div class="block">')   # zero-width range


def test_viz_scatter_trend_line_is_descriptive_only():
    """The OLS overlay must be arithmetically right and must not appear when it would
    be meaningless — a vertical cloud has no slope, and drawing one would invent a
    relationship the data does not contain."""
    import re
    fit = viz.scatter_chart([1, 2, 3, 4], [2, 4, 6, 8], trend_line=True)
    m = re.search(r"slope (-?[\d,.e+-]*?), r (-?[\d.]+)", fit)
    assert round(float(m.group(1).replace(",", "")), 6) == 2.0
    assert round(float(m.group(2)), 6) == 1.0
    assert "not a forecast" in fit
    inverse = viz.scatter_chart([1, 2, 3], [3, 2, 1], trend_line=True)
    assert round(float(re.search(r", r (-?[\d.]+)", inverse).group(1)), 6) == -1.0

    assert "slope" not in viz.scatter_chart([5, 5, 5], [1, 2, 3], trend_line=True)
    assert "slope" not in viz.scatter_chart([1, 2], [1, 2])        # off by default
    # unpaired / unparseable observations are skipped and counted, never plotted at 0
    assert "skipped" in viz.scatter_chart([1, "x", 3], [2, 4, "n/a"])
    assert "skipped" in viz.scatter_chart([1, 2, 3], [1, 2])
    assert "No plottable pairs" in viz.scatter_chart([], [])


def test_viz_stacked_bar_negatives_and_pivot_passthrough():
    """A negative segment must stack *below* the zero line: folding a credit note into
    the positive stack would inflate the bar it reduces. Also pins the documented
    ability to hand `pivot()` output straight to the chart."""
    import re
    import analyse
    h = viz.stacked_bar({"Rev": [("Q1", 10)], "Credit": [("Q1", -4)]})
    rects = re.findall(r'<rect x="[\d.]+" y="([\d.]+)" width="[\d.]+" height="([\d.]+)"', h)
    assert len(rects) == 2
    # Relative order alone is not enough — it holds even if both segments stack
    # upward. Measure against the zero line: one segment must sit wholly above it
    # and the other wholly below.
    zero_y = float(re.search(r'<line x1="[\d.]+" y1="([\d.]+)"[^>]*stroke-width="1.5"',
                             h).group(1))
    tops = [float(y) for y, _ in rects]
    bottoms = [float(y) + float(hh) for y, hh in rects]
    assert min(tops) < zero_y - 1, "positive segment should rise above the zero line"
    assert max(bottoms) > zero_y + 1, "negative segment should drop below the zero line"

    pv = analyse.pivot(["Q", "Seg", "Amt"],
                       [["Q1", "A", "100"], ["Q1", "B", "50"], ["Q2", "A", "120"]],
                       "Q", "Seg", value="Amt")
    chart = viz.stacked_bar(pv, title="From pivot")
    vals = sorted(float(v.replace(",", "")) for v in re.findall(r": ([\d,.-]+)</title>", chart))
    assert vals == [50.0, 100.0, 120.0]
    assert "From pivot" in chart and 'class="legend"' in chart
    assert "No data" in viz.stacked_bar({})


def test_viz_num_reuses_the_engine_parser():
    """The charts must read values exactly as the rest of the toolkit does. A second,
    divergent numeric dialect in viz.py would put a different number on the chart than
    in the working paper."""
    assert viz._num("15%") == 0.15
    assert viz._num("1.2m") == 1200000.0
    assert viz._num("(500)") == -500.0
    assert viz._num("1,234.50") == 1234.5
    assert viz._num("abc") is None and viz._num("") is None
    assert viz._num(True) is None                 # bools are not quantities
    assert viz._num(float("nan")) is None and viz._num(float("inf")) is None


def test_analyse_filter_rows_comparison_semantics():
    """filter_rows standardises hand-rolled filtering, so its comparison rules are the
    whole point. Two regressions are pinned here because both produce a *plausible
    wrong answer* rather than an error:

    - `parse_number('15/02/2026')` returns 15022026 (it strips the separators), so a
      date column compared numerically puts 15 Feb after 1 Mar.
    - falling back to a string compare on a type mismatch makes `'n/a' > 1000` true.
    """
    import analyse
    hdr = ["Customer", "Status", "Amount", "Due", "Note"]
    rows = [
        ["Acme", "Open", "1,200", "01/03/2026", "urgent"],
        ["Borex", "Closed", "900", "15/02/2026", ""],
        ["Credix", "Open", "(500)", "20/04/2026", "credit note"],
        ["Delta", "open", "30", "10/01/2026", "  "],
        ["Echo", "Hold", "n/a", "", "review"],
    ]

    def ids(rs):
        return [r[0] for r in rs]

    # numeric, not lexical: '900' must not beat '1,200'
    out, _ = analyse.filter_rows(hdr, rows, [{"column": "Amount", "op": ">", "value": 1000}])
    assert ids(out) == ["Acme"]
    out, _ = analyse.filter_rows(hdr, rows, [{"column": "Amount", "op": "<", "value": 0}])
    assert ids(out) == ["Credix"]                       # accounting negative

    # dates compare chronologically, not as stripped integers
    out, _ = analyse.filter_rows(hdr, rows, [{"column": "Due", "op": ">",
                                              "value": "01/03/2026"}])
    assert ids(out) == ["Credix"]
    out, _ = analyse.filter_rows(hdr, rows, [{"column": "Due", "op": "between",
                                              "value": ["01/01/2026", "28/02/2026"]}])
    assert ids(out) == ["Borex", "Delta"]

    # between is inclusive of BOTH ends (finance quotes ranges inclusively)
    out, _ = analyse.filter_rows(hdr, rows, [{"column": "Amount", "op": "between",
                                              "value": [30, 900]}])
    assert ids(out) == ["Borex", "Delta"]
    out, _ = analyse.filter_rows(hdr, rows, [{"column": "Amount", "op": "not_between",
                                              "value": [30, 900]}])
    assert ids(out) == ["Acme", "Credix"]

    # a text cell against a numeric threshold is incomparable — dropped AND counted
    out, rep = analyse.filter_rows(hdr, rows, [{"column": "Amount", "op": ">", "value": 0}])
    assert ids(out) == ["Acme", "Borex", "Delta"]
    assert rep["incomparable"] == 1

    # strings are case-insensitive across ==, in, contains
    assert ids(analyse.filter_rows(hdr, rows, [{"column": "Status", "op": "==",
                                                "value": "OPEN"}])[0]) == \
        ["Acme", "Credix", "Delta"]
    assert ids(analyse.filter_rows(hdr, rows, [{"column": "Status", "op": "in",
                                                "value": ["open", "hold"]}])[0]) == \
        ["Acme", "Credix", "Delta", "Echo"]
    assert ids(analyse.filter_rows(hdr, rows, [{"column": "Note", "op": "contains",
                                                "value": "CREDIT"}])[0]) == ["Credix"]

    # whitespace-only is empty
    assert ids(analyse.filter_rows(hdr, rows, [{"column": "Note",
                                                "op": "is_empty"}])[0]) == ["Borex", "Delta"]

    # AND across filters, with per-filter attribution and honest totals
    out, rep = analyse.filter_rows(hdr, rows, [
        {"column": "Status", "op": "==", "value": "open"},
        {"column": "Amount", "op": ">", "value": 100},
    ])
    assert ids(out) == ["Acme"]
    assert (rep["n_in"], rep["n_out"], rep["removed"]) == (5, 1, 4)
    assert [f["removed"] for f in rep["filters"]] == [2, 2]

    assert len(analyse.filter_rows(hdr, rows, [])[0]) == 5      # no filters = passthrough
    assert len(rows) == 5                                       # input not mutated

    # a typo must raise, not quietly match nothing
    for bad, exc in [([{"column": "Nope", "op": "==", "value": 1}], KeyError),
                     ([{"column": "Status", "op": "~=", "value": 1}], ValueError),
                     ([{"column": "Amount", "op": "between", "value": [1]}], ValueError),
                     ([{"column": "Status", "op": "=="}], ValueError)]:
        try:
            analyse.filter_rows(hdr, rows, bad)
            raise AssertionError(f"expected {exc.__name__} for {bad}")
        except exc:
            pass


def test_viz_heatmap_sparkline_waterfall_and_analysis_handoff():
    hm = viz.heatmap([[1, -1], [0.5, 0.2]], row_labels=["A", "B"], col_labels=["X", "Y"],
                     title="Corr", scale="diverging", mid=0)
    assert 'class="chart heatmap"' in hm and "Corr" in hm and "A" in hm
    sp = viz.sparkline([("W1", 10), ("W2", 14), ("W3", 9)], title="Shape")
    assert "spark" in sp and "overall" in sp
    wf = viz.waterfall([
        {"label": "Open", "value": 100, "kind": "start"},
        {"label": "Win", "value": 20, "kind": "delta"},
        {"label": "Loss", "value": -5, "kind": "delta"},
        {"label": "Close", "value": 115, "kind": "total"},
    ], title="Bridge")
    assert "Bridge" in wf and wf.count("<rect") >= 4
    empty = viz.heatmap([], title="Empty")
    assert "No data" in empty

    analysis = {
        "results": [
            {"op": "breakdown", "name": "By customer", "result": {
                "by": "Customer",
                "groups": [
                    {"key": "Acme", "count": 2, "total": "100", "share": "0.625"},
                    {"key": "Beta", "count": 1, "total": "60", "share": "0.375"},
                ],
            }},
            {"op": "period_series", "name": "Monthly", "result": {
                "grain": "month",
                "periods": [
                    {"period": "2026-01", "count": 1, "total": "100", "delta": None},
                    {"period": "2026-02", "count": 1, "total": "130", "delta": "30"},
                ],
            }},
            {"op": "correlation_matrix", "name": "Corr", "result": {
                "columns": ["A", "B"], "matrix": [[1.0, 0.5], [0.5, 1.0]],
            }},
            {"op": "forecast", "name": "skip me", "result": {"value": 1}},
        ]
    }
    specs = viz.suggest_blocks_from_analysis(analysis)
    assert [s["type"] for s in specs] == ["section", "section", "section"]
    assert any(b["type"] == "waterfall" for s in specs for b in s["blocks"])
    assert any(b["type"] == "heatmap" and b.get("scale") == "diverging"
               for s in specs for b in s["blocks"])
    filtered = viz.suggest_blocks_from_analysis(analysis, ops=["breakdown"])
    assert len(filtered) == 1 and filtered[0]["title"] == "By customer"
    html_blocks = viz.blocks_from_analysis(analysis)
    page = viz.dashboard("Insight board", html_blocks, as_of="18 Jul 2026")
    assert "By customer" in page and "period bridge" in page and "heatmap" in page


# --------------------------------------------------------------------------- #
# 7 & 8. form extraction: next-line + dotted-leader layouts
# --------------------------------------------------------------------------- #
_FORM_FIELDS = [
    {"name": "Investor", "labels": ["investor"], "type": "text"},
    {"name": "Settlement bank", "labels": ["settlement bank", "bank"], "type": "text"},
    {"name": "Fee", "labels": ["fee"], "type": "currency", "currency": "GBP"},
]
_FORM_TEXT = ("Subscription confirmation\n"
              "Investor: Acme Pension Fund\n"
              "Settlement bank\n"          # label alone -> value on the next line
              "HSBC London\n"
              "Fee .......... GBP 2,500\n")  # dotted leader


def test_next_line_form_field():
    rec, _ = extract.extract_fields(None, _FORM_FIELDS, text=_FORM_TEXT)
    assert rec["Settlement bank"] == "HSBC London"


def test_dotted_leader_form_field():
    rec, _ = extract.extract_fields(None, _FORM_FIELDS, text=_FORM_TEXT)
    assert rec["Fee"] == Decimal("2500")


def test_next_line_stops_at_following_label():
    # if a field's value is genuinely absent, the next-line search must not grab the next label
    fields = [{"name": "A", "labels": ["field a"], "type": "text"},
              {"name": "B", "labels": ["field b"], "type": "text"}]
    rec, flags = extract.extract_fields(None, fields, text="Field A\nField B\nvalue-b\n")
    assert rec["A"] == "" and rec["B"] == "value-b"
    assert any(f["field"] == "A" for f in flags)


# --------------------------------------------------------------------------- #
# 9. currency code_target splits the code into its own column
# --------------------------------------------------------------------------- #
def test_code_target_currency_split():
    raw = [["Amount"], ["S$ 1,000"], ["US$ 2,000"], ["$ 3,000"]]
    recipe = {"columns": [{"source": "Amount", "target": "Amount", "type": "currency",
                           "currency": "SGD", "code_target": "Currency"}]}
    header, rows, _ = dataclean.apply_recipe(raw, recipe)
    assert header == ["Amount", "Currency"]
    assert [r[1] for r in rows] == ["SGD", "USD", "SGD"]    # bare $ resolved to expected SGD
    assert rows[0][0] == Decimal("1000")


# --------------------------------------------------------------------------- #
# 10. categorical value-map proposal
# --------------------------------------------------------------------------- #
def test_categorical_value_map_proposal():
    clusters = dataclean.propose_value_map(["USA", "U.S.A.", "USA", "usa",
                                            "United States", "Canada", "canada"])
    top = clusters[0]
    assert top["canonical"] == "USA" and set(top["variants"]) >= {"U.S.A.", "usa"}
    vmap = dataclean.value_map_from_clusters(clusters)
    val, note, _ = dataclean._convert("u.s.a.", {"type": "categorical", "value_map": vmap})
    assert val == "USA" and "standardised" in note
    # master list snaps to the master spelling
    cm = dataclean.propose_value_map(["united states", "UNITED STATES"], master=["United States"])
    assert cm[0]["canonical"] == "United States" and cm[0]["from_master"]


# --------------------------------------------------------------------------- #
# 11. PDF table extraction — engine scoring (deterministic) + guarded smoke test
# --------------------------------------------------------------------------- #
def test_pdf_table_scoring_and_selection():
    good = [["A", "B"], ["1", "2"], ["3", "4"]]            # 3x2 consistent
    onecol = [["just text"], ["more text"]]                # not table-shaped
    assert ingest._table_score(good) > 0
    assert ingest._table_score(onecol) == 0
    # better-scoring engine wins; ties go to pdfplumber (preferred for messy tables)
    assert ingest._choose_tables([good], [onecol]) == ([good], "pdfplumber")
    assert ingest._choose_tables([onecol], [good]) == ([good], "pymupdf")
    assert ingest._choose_tables([good], [good]) == ([good], "pdfplumber")
    assert ingest._choose_tables([], []) == ([], None)


def test_pdf_read_smoke():
    try:
        import fitz  # noqa: F401
    except ImportError:
        print("  [skip] PyMuPDF not installed — PDF smoke test skipped")
        return
    import fitz
    d = tempfile.mkdtemp()
    pdf = Path(d) / "t.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 100), "Investor    Commitment    Currency")
    page.insert_text((72, 120), "Acme        1000           USD")
    page.insert_text((72, 140), "Beta        2000           SGD")
    doc.save(str(pdf))
    doc.close()
    rows, note = ingest.read_pdf(str(pdf))
    flat = " ".join(str(c) for r in rows for c in r)
    assert "Acme" in flat and "1000" in flat, (note, rows)


def test_amount_date_missing_dates_do_not_match():
    A = [{"amount": "100.00", "date": ""}]
    B = [{"amount": "100.00", "date": "01 Jun 2026"}]
    res = reconcile.match(A, B, amount="amount", date="date", mode="amount_date")
    assert not res["matched"] and not res["ambiguous"]
    assert len(res["a_only"]) == 1 and len(res["b_only"]) == 1

    bad = reconcile.match([{"amount": "100.00", "date": "not a date"}], B,
                          amount="amount", date="date", mode="amount_date")
    assert not bad["matched"] and len(bad["a_only"]) == 1 and len(bad["b_only"]) == 1


def test_key_mode_missing_amount_flagged():
    res = reconcile.match([{"ref": "R1", "amount": ""}],
                          [{"ref": "R1", "amount": "0.00"}],
                          key="ref", amount="amount")
    assert not res["matched"] and not res["value_diffs"]
    assert len(res["parse_errors"]) == 1
    assert any(e["category"] == "parse_error" for e in reconcile.triage(res))


def test_aggregation_currency_mismatch_rejected():
    one = {
        "matched": [], "value_diffs": [], "currency_diffs": [], "currency_unknown": [],
        "parse_errors": [], "dup_a": [], "dup_b": [], "ambiguous": [],
        "a_only": [{"a": {"row": {"batch": "B1"}, "amt": Decimal("100"), "dt": None, "ccy": "USD"}}],
        "b_only": [{"b": {"row": {"batch": "B1"}, "amt": Decimal("40"), "dt": None, "ccy": "USD"}},
                   {"b": {"row": {"batch": "B1"}, "amt": Decimal("60"), "dt": None, "ccy": "SGD"}}],
    }
    assert reconcile.propose_aggregations(one, group_col="batch") == []

    many = {
        "matched": [], "value_diffs": [], "currency_diffs": [], "currency_unknown": [],
        "parse_errors": [], "dup_a": [], "dup_b": [], "ambiguous": [],
        "a_only": [{"a": {"row": {"batch": "B2"}, "amt": Decimal("100"), "dt": None, "ccy": "USD"}}],
        "b_only": [{"b": {"row": {"batch": "B2"}, "amt": Decimal("100"), "dt": None, "ccy": "SGD"}},
                   {"b": {"row": {"batch": "B2"}, "amt": Decimal("0"), "dt": None, "ccy": "SGD"}}],
    }
    assert reconcile.propose_aggregations(many, group_col="batch") == []


def test_apply_recipe_empty_input():
    header, rows, log = dataclean.apply_recipe([], {"columns": []})
    assert header == [] and rows == []
    assert log["rows_out"] == 0 and "empty input" in log["message"]


def test_col_index_ambiguous_short_name():
    assert dataclean._col_index(["Invoice date", "Payment date"], "date") is None
    assert dataclean._col_index(["Amount", "Payment date"], "date") == 1


def test_bool_unrecognised_kept_raw():
    raw = [["Active"], ["maybe"], ["N/A"], ["yes"], ["0"]]
    recipe = {"columns": [{"source": "Active", "target": "Active", "type": "bool"}]}
    _, rows, log = dataclean.apply_recipe(raw, recipe)
    assert rows == [["maybe"], ["N/A"], [True], [False]]
    assert [f["value"] for f in log["flagged"]] == ["maybe", "N/A"]


def test_md_escape_pipe_in_value():
    log = {"header": 0, "dropped": {"blank": 0, "totals": 0}, "transforms": [],
           "flagged": [{"row": 1, "column": "A|B", "value": "x|y", "reason": "bad|value"}],
           "duplicates": [], "validation": [], "rows_in": 1, "rows_out": 1}
    report = dataclean.render_report(log)
    assert "A\\|B" in report and "x\\|y" in report and "bad\\|value" in report

    rreport = reconcile.render_report(
        {"rag": "AMBER", "pct_reconciled": 0, "matched": 0, "aggregated": 0,
         "exceptions": 1, "material_or_escalate": 0, "value_matched": Decimal("0"),
         "value_in_exception": Decimal("1"), "by_category": {"x|y": {"n": 1, "val": Decimal("1")}}},
        [{"category": "x|y", "magnitude": Decimal("1"), "materiality": "immaterial",
          "probable_cause": "a|b", "suggested_action": "c|d"}])
    assert "x\\|y" in rreport and "a\\|b" in rreport


def test_currency_code_case_insensitive():
    assert dataclean.parse_currency("sgd 100")[0] == dataclean.parse_currency("SGD 100")[0]
    assert dataclean.parse_currency("sgd 100")[0][1] == "SGD"
    assert reconcile.to_currency("sgd 100") == "SGD"


def test_european_number_format_flagged():
    n, note = dataclean.parse_number("1.234,56")
    assert n is None and "European" in note
    res, cnote = dataclean.parse_currency("EUR 1.234,56")
    assert res is None and "European" in cnote


def test_get_table_engine_param():
    old = ingest.extract_pdf_table
    calls = []

    def fake(path, page, index=0, engine=None):
        calls.append((path, page, index, engine))
        return [["ok"]]

    ingest.extract_pdf_table = fake
    try:
        assert extract.get_table("sample.pdf", page=2, index=3, engine="pdfplumber") == [["ok"]]
        assert calls == [("sample.pdf", 2, 3, "pdfplumber")]
    finally:
        ingest.extract_pdf_table = old


def test_logo_restricted_to_images():
    d = tempfile.mkdtemp()
    txt = Path(d) / "logo.txt"
    txt.write_text("not image", encoding="utf-8")
    html = viz._logo_for({"logo_path": str(txt), "brand_name": "Acme"})
    assert "<img" not in html and "Acme" in html


def test_pii_egress_nric_pattern():
    hook = _ROOT / "hooks" / "pii_egress_guard.py"
    payload = {"tool_input": {"query": "Can you search S1234567D payment status?"}}
    proc = subprocess.run([sys.executable, str(hook)], input=json.dumps(payload),
                          text=True, capture_output=True, check=False)
    assert proc.returncode == 0
    out = json.loads(proc.stdout)
    assert out["hookSpecificOutput"]["permissionDecision"] == "ask"


# --------------------------------------------------------------------------- #
# Day-to-day finance strengthening (v0.2.1): debit/credit columns, case-insensitive
# column names, flip_b, balance completeness, ageing, tax hint, per-currency summary
# --------------------------------------------------------------------------- #
def test_column_names_resolved_case_insensitively():
    # a bank CSV says "Amount" / "Date"; the config says amount/date — must still match
    A = [{"Amount": "100.00", "Date": "01/06/2026"}]
    B = [{"amount": "100.00", "date": "01/06/2026"}]
    res = reconcile.match(A, B, amount="amount", date="date", mode="amount_date")
    assert len(res["matched"]) == 1, res


def test_debit_credit_columns_make_signed_amount():
    # bank export splits Debit/Credit; ledger holds one signed amount — they must reconcile
    bank = [{"Date": "01/06/2026", "Debit": "", "Credit": "1,000.00"},
            {"Date": "02/06/2026", "Debit": "250.00", "Credit": ""}]
    ledger = [{"date": "01/06/2026", "amount": "-1,000.00"},
              {"date": "02/06/2026", "amount": "250.00"}]
    res = reconcile.match(bank, ledger, amount="amount", debit="debit", credit="credit",
                          date="date", mode="amount_date")
    assert len(res["matched"]) == 2, res


def test_flip_b_reconciles_opposite_sign_conventions():
    A = [{"k": "1", "amount": "100.00"}]
    B = [{"k": "1", "amount": "-100.00"}]
    assert not reconcile.match(A, B, key="k", amount="amount")["matched"]
    res = reconcile.match(A, B, key="k", amount="amount", flip_b=True)
    assert len(res["matched"]) == 1, res


def test_balance_completeness_check():
    rows = [{"amount": "1,000.00"}, {"amount": "-250.00"}]
    ok = reconcile.check_balance(rows, amount="amount", opening="100.00", closing="850.00")
    assert ok["ties"] is True and ok["movement"] == Decimal("750.00")
    # a truncated extract must NOT tie
    bad = reconcile.check_balance(rows[:1], amount="amount", opening="100.00", closing="850.00")
    assert bad["ties"] is False
    assert "DOES NOT TIE" in reconcile.render_balance_check(bad, "Bank")


def test_aging_of_one_sided_items():
    res = reconcile.match([{"amount": "500.00", "date": "01 Jun 2026"}], [],
                          amount="amount", date="date", mode="amount_date")
    exc = reconcile.triage(res, as_of="30 Jun 2026")
    assert exc[0]["age_days"] == 29
    assert "aged 29d" in exc[0]["probable_cause"]
    # no as_of -> no ageing (deterministic default)
    assert reconcile.triage(res)[0]["age_days"] is None


def test_tax_hint_on_net_vs_gross_mismatch():
    res = reconcile.match([{"k": "1", "amount": "1000.00"}],
                          [{"k": "1", "amount": "1090.00"}], key="k", amount="amount")
    exc = reconcile.triage(res)
    assert exc[0]["category"] == "amount_mismatch"
    assert "GST/VAT/WHT" in exc[0]["probable_cause"]
    # a non-tax-shaped difference gets no hint
    res2 = reconcile.match([{"k": "1", "amount": "1000.00"}],
                           [{"k": "1", "amount": "1234.00"}], key="k", amount="amount")
    assert "GST/VAT/WHT" not in reconcile.triage(res2)[0]["probable_cause"]


def test_summary_per_currency_breakdown():
    A = [{"ref": "R1", "amount": "100.00", "ccy": "USD"},
         {"ref": "R2", "amount": "50.00", "ccy": "SGD"},
         {"ref": "R3", "amount": "70.00", "ccy": "SGD"}]
    B = [{"ref": "R1", "amount": "100.00", "ccy": "USD"},
         {"ref": "R2", "amount": "50.00", "ccy": "SGD"}]
    res = reconcile.match(A, B, key="ref", amount="amount", currency="ccy")
    exc = reconcile.triage(res)
    s = reconcile.summarise(res, exc)
    assert s["by_currency"]["USD"]["matched_val"] == Decimal("100.00")
    assert s["by_currency"]["SGD"]["matched_val"] == Decimal("50.00")
    assert s["by_currency"]["SGD"]["exception_val"] == Decimal("70.00")
    report = reconcile.render_report(s, exc)
    assert "Mixed currencies" in report


# --------------------------------------------------------------------------- #
# 16. Large-file streaming (scripts/streaming.py + ingest.read_large)
# --------------------------------------------------------------------------- #
def test_streaming_count_rows():
    """Row counts via read_only iter_rows; strategy gate for a small sheet is direct."""
    import streaming
    d = tempfile.mkdtemp()
    path = Path(d) / "rows.xlsx"
    rows = [["id", "amount", "ccy"]] + [[i, i * 1.5, "USD"] for i in range(1, 501)]
    _write_xlsx(path, [("Ledger", rows), ("Empty", None)])
    counts = streaming.count_rows(str(path))
    assert counts["sheets"]["Ledger"] == 501          # header + 500 data
    assert counts["sheets"]["Empty"] == 0
    assert counts["total"] == 501
    one = streaming.count_rows(str(path), sheet="Ledger")
    assert one["sheets"] == {"Ledger": 501}
    assert streaming.choose_strategy(str(path), sheet="Ledger") == "direct"


def test_streaming_excel_to_parquet_and_read_large():
    import streaming
    if not streaming.pyarrow_available():
        print("  SKIP  test_streaming_excel_to_parquet_and_read_large (no pyarrow/pandas)")
        return
    import pandas as pd
    d = tempfile.mkdtemp()
    path = Path(d) / "stream_me.xlsx"
    n = 2500
    rows = [["txn_id", "amount", "region"]] + [
        [f"T{i}", float(i), "APAC" if i % 2 == 0 else "EMEA"] for i in range(n)
    ]
    _write_xlsx(path, [("Data", rows)])
    pq = Path(d) / "out.parquet"
    written = streaming.stream_excel_to_parquet(str(path), str(pq), sheet="Data", chunk_size=800)
    assert written == n
    df = pd.read_parquet(pq)
    assert len(df) == n
    assert list(df.columns) == ["txn_id", "amount", "region"]

    # Force stream strategy via threshold monkeypatch for ingest.read_large
    old_d, old_p = streaming.DIRECT_THRESHOLD, streaming.PARQUET_CACHE_THRESHOLD
    streaming.DIRECT_THRESHOLD = 10
    streaming.PARQUET_CACHE_THRESHOLD = 100
    try:
        assert streaming.choose_strategy(str(path), sheet="Data") == "stream"
        frame, note = ingest.read_large(str(path), sheet="Data", cache_dir=d)
        assert "stream" in note
        assert len(frame) == n
    finally:
        streaming.DIRECT_THRESHOLD = old_d
        streaming.PARQUET_CACHE_THRESHOLD = old_p


def test_parquet_cache_preserves_none_and_nan_literals():
    """parquet_cache must not null real string values 'None' / 'nan' (review M1)."""
    import streaming
    if not streaming.pyarrow_available():
        print("  SKIP  test_parquet_cache_preserves_none_and_nan_literals (no pyarrow/pandas)")
        return
    d = tempfile.mkdtemp()
    path = Path(d) / "names.xlsx"
    # 12 data rows → force parquet_cache with lowered thresholds
    rows = [["customer", "note"]] + [
        ["None Corp", "nan handling pending"],
        ["Acme", "ok"],
    ] * 6
    _write_xlsx(path, [("Data", rows)])
    old_d, old_p = streaming.DIRECT_THRESHOLD, streaming.PARQUET_CACHE_THRESHOLD
    streaming.DIRECT_THRESHOLD = 2
    streaming.PARQUET_CACHE_THRESHOLD = 10_000
    try:
        assert streaming.choose_strategy(str(path), sheet="Data") == "parquet_cache"
        frame, note = ingest.read_large(str(path), sheet="Data", cache_dir=d)
        assert "parquet_cache" in note
        assert "None Corp" in set(frame["customer"].astype(str))
        assert "nan handling pending" in set(frame["note"].astype(str))
        # Second read hits cache — still preserved
        frame2, note2 = ingest.read_large(str(path), sheet="Data", cache_dir=d)
        assert "cache hit" in note2
        assert "None Corp" in set(frame2["customer"].astype(str))
    finally:
        streaming.DIRECT_THRESHOLD = old_d
        streaming.PARQUET_CACHE_THRESHOLD = old_p


def test_optimize_dtypes_saves_memory():
    import streaming
    try:
        import pandas as pd
        import numpy as np
    except ImportError:
        print("  SKIP  test_optimize_dtypes_saves_memory (no pandas/numpy)")
        return
    n = 50_000
    # Mixed finance-shaped frame: wide int64/float64 + long repeated object labels.
    # Category conversion on long strings is where most of the saving comes from.
    regions = np.array(
        ["Asia-Pacific regional coverage office"] * (n // 4)
        + ["Europe Middle East Africa desk"] * (n // 4)
        + ["Americas institutional coverage"] * (n // 4)
        + ["Global multi-strategy allocation"] * (n // 4),
        dtype=object,
    )
    statuses = np.array(
        ["open - pending settlement review"] * (n // 2)
        + ["closed - archived prior period"] * (n // 2),
        dtype=object,
    )
    df = pd.DataFrame({
        "id": np.arange(n, dtype="int64"),
        "flag": np.zeros(n, dtype="int64"),
        "amt": np.linspace(0.0, 1.0, n).astype("float64"),
        "region": regions,
        "status": statuses,
    })
    before = df.memory_usage(deep=True).sum()
    out = streaming.optimize_dtypes(df)
    after = out.memory_usage(deep=True).sum()
    saved = (1 - after / before) * 100
    assert saved >= 40, f"expected ≥40% memory save, got {saved:.1f}%"
    assert str(out["region"].dtype) == "category"
    assert str(out["status"].dtype) == "category"


# --------------------------------------------------------------------------- #
# 17. Image / chart extraction (skills/data-extract/scripts/image_extract.py)
# --------------------------------------------------------------------------- #
def _load_image_extract():
    import importlib.util
    mod_path = _ROOT / "skills" / "data-extract" / "scripts" / "image_extract.py"
    spec = importlib.util.spec_from_file_location("image_extract", mod_path)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


def test_parse_markdown_table_numeric_cleanup():
    ie = _load_image_extract()
    text = """
Here is the table:

| Region | Revenue | Growth | Fee |
|---|---|---|---|
| APAC | $1,234.50 | 12.5% | £90 |
| EMEA | US$2,000 | 8% | €1,100 |
| AMER | (500) | -3.2% | S$250 |

Done.
"""
    df = ie.parse_markdown_table(text)
    assert df is not None
    assert list(df.columns) == ["Region", "Revenue", "Growth", "Fee"]
    assert len(df) == 3
    # Works with pandas DataFrame or the no-pandas stand-in
    row0 = df.iloc[0] if hasattr(df, "iloc") else dict(zip(df.columns, df._data[0]))
    row1 = df.iloc[1] if hasattr(df, "iloc") else dict(zip(df.columns, df._data[1]))
    row2 = df.iloc[2] if hasattr(df, "iloc") else dict(zip(df.columns, df._data[2]))
    assert row0["Revenue"] == 1234.5
    assert abs(row0["Growth"] - 0.125) < 1e-9
    assert row1["Revenue"] == 2000
    assert row2["Revenue"] == -500
    assert row1["Fee"] == 1100


def test_image_extract_chart_and_table_mocked():
    ie = _load_image_extract()
    try:
        from PIL import Image
    except ImportError:
        print("  SKIP  test_image_extract_chart_and_table_mocked (no Pillow)")
        return

    d = Path(tempfile.mkdtemp())
    chart = d / "sales_bar_chart.png"
    table = d / "ledger_table.png"
    Image.new("RGB", (320, 200), color=(40, 80, 160)).save(chart)
    Image.new("RGB", (400, 240), color=(240, 240, 240)).save(table)

    chart_md = """
| Category | Value |
|---|---|
| Q1 | 120 |
| Q2 | 150 |
| Q3 | 90 |
| Q4 | 200 |
"""
    table_md = """
| Investor | Commit | Close |
|---|---|---|
| Acme Pension | 1000000 | 12 Jun 2026 |
| Beta FO | 2500000 | 13 Jun 2026 |
| Gamma LP | 750000 | 14 Jun 2026 |
"""

    class _Resp:
        def __init__(self, text, status=200):
            self.status_code = status
            self.text = text
            self._text = text

        def json(self):
            return {
                "choices": [{"message": {"content": self._text}}],
                "usage": {"prompt_tokens": 10, "completion_tokens": 20},
                "model": "mock-vision",
            }

    calls = {"n": 0}

    def fake_post(url, headers=None, json=None, timeout=None):
        calls["n"] += 1
        # Route by which prompt was used
        prompt = json["messages"][0]["content"][0]["text"]
        if "chart title" in prompt.lower() or "data point" in prompt.lower():
            return _Resp(chart_md)
        return _Resp(table_md)

    cache = d / "cache"
    r1 = ie.extract_image(
        str(chart), api_key="test-key", model="mock-vision",
        cache_dir=cache, _request=fake_post,
    )
    assert r1.get("error") is None, f"chart extract error: {r1.get('error')}"
    assert r1["cached"] is False
    assert r1["type"] == "chart"
    assert r1["dataframe"] is not None and len(r1["dataframe"]) == 4
    cats = list(r1["dataframe"]["Category"])
    assert set(cats) == {"Q1", "Q2", "Q3", "Q4"}
    vals = [float(v) for v in r1["dataframe"]["Value"]]
    assert sum(vals) == 560

    # Cache hit on second call
    r1b = ie.extract_image(
        str(chart), api_key="test-key", model="mock-vision",
        cache_dir=cache, _request=fake_post,
    )
    assert r1b["cached"] is True
    assert calls["n"] == 1

    r2 = ie.extract_image(
        str(table), api_key="test-key", model="mock-vision",
        cache_dir=cache, _request=fake_post,
    )
    assert r2.get("error") is None, f"table extract error: {r2.get('error')}"
    assert r2["type"] == "table"
    assert r2["dataframe"] is not None
    assert list(r2["dataframe"].columns) == ["Investor", "Commit", "Close"]
    assert len(r2["dataframe"]) == 3
    assert int(r2["dataframe"].iloc[0]["Commit"]) == 1_000_000


def test_image_extract_batch_and_compress():
    ie = _load_image_extract()
    try:
        from PIL import Image
    except ImportError:
        print("  SKIP  test_image_extract_batch_and_compress (no Pillow)")
        return

    d = Path(tempfile.mkdtemp())
    imgs = d / "shots"
    imgs.mkdir()
    # Five small images + one oversized to exercise compression
    for i in range(5):
        Image.new("RGB", (64, 64), color=(i * 40, 100, 120)).save(imgs / f"table_{i}.png")
    big = imgs / "huge_photo.png"
    # Oversized on the long edge — compression path must resize (byte shrink is not
    # guaranteed for already-tiny solid-colour PNGs).
    Image.new("RGB", (2200, 1600), color=(200, 100, 50)).save(big, compress_level=1)

    md = "| A | B |\n|---|---|\n| 1 | 2 |\n| 3 | 4 |\n"

    class _Resp:
        status_code = 200
        text = md

        def json(self):
            return {
                "choices": [{"message": {"content": md}}],
                "usage": {},
                "model": "mock-vision",
            }

    def fake_post(url, headers=None, json=None, timeout=None):
        return _Resp()

    out = d / "batch.xlsx"
    summary = ie.extract_batch(
        str(imgs), str(out),
        api_key="test-key", model="mock-vision",
        cache_dir=d / "cache", _request=fake_post,
    )
    assert out.is_file(), "batch xlsx was not written"
    assert summary["count"] == 6, f"expected 6 images, got {summary['count']}"
    assert "combined" in summary["sheets"]

    import io
    compressed_bytes, mime, was_compressed = ie.compress_image(
        big, max_bytes=50_000, max_px=1024,
    )
    assert was_compressed is True
    out_img = Image.open(io.BytesIO(compressed_bytes))
    assert max(out_img.size) <= 1024
    assert max(Image.open(big).size) > 1024
# 17. PowerPoint ingest (read_pptx) + CJK / i18n fonts (viz)
# --------------------------------------------------------------------------- #
def _write_sample_pptx(path, *, with_image_only=True):
    """Build a multi-slide deck: tables on slides 1–2, optional image-only slide 3."""
    from pptx import Presentation
    from pptx.util import Inches
    import io
    prs = Presentation()
    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    box = s1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(0.4))
    box.text_frame.text = "Q2 commitments"
    t1 = s1.shapes.add_table(3, 2, Inches(0.5), Inches(0.8), Inches(4), Inches(1.4)).table
    t1.cell(0, 0).text, t1.cell(0, 1).text = "Investor", "Commit"
    t1.cell(1, 0).text, t1.cell(1, 1).text = "Acme", "1000000"
    t1.cell(2, 0).text, t1.cell(2, 1).text = "Beta", "2500000"
    # Second table on the same slide (multi-table)
    t1b = s1.shapes.add_table(2, 2, Inches(0.5), Inches(2.5), Inches(3), Inches(0.9)).table
    t1b.cell(0, 0).text, t1b.cell(0, 1).text = "Region", "Share"
    t1b.cell(1, 0).text, t1b.cell(1, 1).text = "APAC", "60%"

    s2 = prs.slides.add_slide(blank)
    t2 = s2.shapes.add_table(2, 2, Inches(0.5), Inches(0.5), Inches(3), Inches(0.9)).table
    t2.cell(0, 0).text, t2.cell(0, 1).text = "Metric", "Value"
    t2.cell(1, 0).text, t2.cell(1, 1).text = "AUM", "3.5bn"

    if with_image_only:
        s3 = prs.slides.add_slide(blank)
        try:
            from PIL import Image
            buf = io.BytesIO()
            Image.new("RGB", (24, 24), color=(200, 40, 40)).save(buf, format="PNG")
            buf.seek(0)
            s3.shapes.add_picture(buf, Inches(1), Inches(1), width=Inches(1))
        except ImportError:
            # Still image-only if we add nothing — empty slide counts as image-only too.
            pass

    prs.save(path)


def test_pptx_multi_slide_tables_and_read_any():
    try:
        import pptx  # noqa: F401
    except ImportError:
        print("  SKIP  test_pptx_multi_slide_tables_and_read_any (no python-pptx)")
        return
    d = tempfile.mkdtemp()
    path = Path(d) / "deck.pptx"
    _write_sample_pptx(path, with_image_only=True)
    rows, note = ingest.read_pptx(str(path))
    # Table rows only — title text frames must NOT pollute the row list (review M1)
    assert any(r == ["Investor", "Commit"] for r in rows)
    assert any(r == ["Acme", "1000000"] for r in rows)
    assert any(r == ["Region", "Share"] for r in rows)
    assert any(r == ["Metric", "Value"] for r in rows)
    assert not any(r == ["Q2 commitments"] for r in rows)
    assert all(len(r) >= 2 for r in rows), "text frames must not create 1-col rows"
    assert "tables on slide(s) 1, 2" in note
    assert "image-only slide(s) 3" in note
    assert "Q2 commitments" in note          # titles live in the note
    # read_any dispatches
    rows2, note2 = ingest.read_any(str(path))
    assert rows2 == rows and "pptx" in note2


def test_pptx_legacy_ppt_rejected_and_image_only_note():
    try:
        import pptx  # noqa: F401
    except ImportError:
        print("  SKIP  test_pptx_legacy_ppt_rejected_and_image_only_note (no python-pptx)")
        return
    try:
        ingest.read_any("legacy.ppt")
        assert False, "expected ValueError for .ppt"
    except ValueError as e:
        assert "Legacy .ppt not supported" in str(e)
        assert "convert to .pptx" in str(e)

    d = tempfile.mkdtemp()
    path = Path(d) / "pics.pptx"
    _write_sample_pptx(path, with_image_only=True)
    _, note = ingest.read_pptx(str(path))
    assert "image-only" in note
    assert "vision-model" in note


def test_has_cjk_detection():
    assert viz._has_cjk("Hello") is False
    assert viz._has_cjk("销售额") is True
    assert viz._has_cjk("Q2 売上") is True          # CJK ideograph + kana
    assert viz._has_cjk("안녕하세요") is True       # Hangul
    assert viz._has_cjk("") is False
    assert "arabic" in viz._detect_scripts("الإيرادات")
    assert "thai" in viz._detect_scripts("ยอดขาย")
    assert viz._needs_i18n_fonts("Hello world 2026") is False


def test_cjk_font_stack_conditional():
    # English-only: theme fonts unchanged (no CJK names injected)
    en = viz.dashboard("Operations dashboard",
                       [viz.bar_chart([("Mon", 4), ("Tue", 7)], title="Tasks by day")],
                       subtitle="Weekly")
    assert "Microsoft YaHei" not in en
    assert "Noto Sans CJK" not in en
    assert 'lang="en"' in en
    assert 'dir="ltr"' in en
    # Default body font still present
    assert "Inter" in en or "Segoe UI" in en

    # Chinese labels: CJK fallback stack applied; SVG chart text inherits via CSS
    zh = viz.dashboard("销售仪表板",
                       [viz.bar_chart([("一月", 4), ("二月", 7)], title="销售额")],
                       subtitle="季度回顾")
    assert "Microsoft YaHei" in zh or "Noto Sans CJK" in zh or "PingFang SC" in zh
    assert "销售额" in zh
    # Chart CSS uses the (augmented) font stack — SVG <text> inherits
    assert ".chart .v{font:11px" in zh.replace(" ", "") or "font:11px" in zh

    # Fully Arabic page → RTL + Arabic font names
    ar = viz.dashboard("لوحة القيادة", [viz.kpi_row([{"label": "الإيرادات", "value": 12}])])
    assert 'dir="rtl"' in ar
    assert "Noto Naskh Arabic" in ar or "Noto Sans Arabic" in ar

    # Mostly English with one Arabic label must NOT flip the whole page (review m2)
    mixed = viz.dashboard(
        "Operations dashboard",
        [viz.kpi_row([{"label": "Open tasks", "value": 12},
                      {"label": "الإيرادات", "value": 3}])],
        subtitle="Weekly",
    )
    assert 'dir="ltr"' in mixed
    assert "Noto Naskh Arabic" in mixed or "Noto Sans Arabic" in mixed


# --------------------------------------------------------------------------- #
# 18. New analysis functions: concentration, pivot, distribution, trend
# --------------------------------------------------------------------------- #
def test_concentration_hhi_and_classification():
    import analyse
    # Highly concentrated: one group = 90% of total
    c = analyse.concentration(["900", "50", "30", "20"])
    assert c["n"] == 4 and c["n_groups"] == 4
    assert c["hhi"] is not None and c["hhi"] > Decimal(5000), c
    assert c["classification"] == "highly concentrated", c
    assert c["top_n_share"] is not None and c["top_n_share"] > Decimal("0.9")
    # Fragmented: 10 equal groups → HHI = 10 * (10%)² = 1000
    c2 = analyse.concentration(["10"] * 10)
    assert c2["hhi"] < Decimal(1500), c2
    assert c2["classification"] == "fragmented", c2
    # Negatives → unreliable
    c3 = analyse.concentration(["100", "-50", "60"])
    assert c3["hhi"] is None and "negatives" in c3["classification"]
    # Empty
    c4 = analyse.concentration([])
    assert c4["hhi"] is None and c4["classification"] == "no data"


def test_pivot_cross_tab_sum_and_count():
    import analyse
    header = ["Date", "Customer", "Region", "Amount"]
    rows = [
        ["05/01/2026", "Alpha", "North", "1,000"],
        ["12/01/2026", "Beta", "South", "2,000"],
        ["20/01/2026", "Alpha", "North", "1,500"],
        ["15/03/2026", "Gamma", "South", "(500)"],
        ["18/03/2026", "Alpha", "North", "9,000"],
        ["25/03/2026", "Delta", "East", "junk"],
        ["", "Echo", "", "250"],
    ]
    pv = analyse.pivot(header, rows, "Region", "Customer", value="Amount")
    assert pv["n_rows"] == 4 and pv["n_cols"] == 5  # incl (blank)
    assert "North" in pv["row_keys"] and "Alpha" in pv["col_keys"]
    ni, ai = pv["row_keys"].index("North"), pv["col_keys"].index("Alpha")
    assert pv["matrix"][ni][ai] == Decimal("11500"), pv["matrix"][ni][ai]
    assert pv["grand_total"] == Decimal("13250"), pv["grand_total"]
    assert pv["skipped"] == 1  # the "junk" amount
    # count aggfunc (no value) — auto-selected
    pvc = analyse.pivot(header, rows, "Region", "Customer")
    assert pvc["aggfunc"] == "count" and pvc["measure"] == "rows"
    assert pvc["matrix"][pvc["row_keys"].index("North")][pvc["col_keys"].index("Alpha")] == Decimal(3)
    # M1 fix: blank amount cells don't crash — they're skipped, not appended as None
    blank_header = ["Region", "Customer", "Amount"]
    blank_rows = [
        ["North", "Alpha", "100"],
        ["North", "Alpha", ""],     # blank amount — must not crash
        ["North", "Beta", "junk"],   # unparseable — skipped, counted
        ["South", "Gamma", "200"],
    ]
    pvb = analyse.pivot(blank_header, blank_rows, "Region", "Customer", value="Amount")
    ni2 = pvb["row_keys"].index("North")
    ai2 = pvb["col_keys"].index("Alpha")
    assert pvb["matrix"][ni2][ai2] == Decimal("100"), pvb["matrix"][ni2][ai2]  # only the 100
    assert pvb["skipped"] == 1  # the "junk"


def test_distribution_skewness_kurtosis():
    import analyse
    # Symmetric: 1,2,3,4,5 → skewness ~0
    d = analyse.distribution(["1", "2", "3", "4", "5"])
    assert d["n"] == 5 and d["skewness"] is not None
    assert abs(d["skewness"]) < 0.1, d["skewness"]
    assert d["classification"] == "symmetric", d
    # Right-skewed with outlier → heavy skew
    d2 = analyse.distribution(["1", "1", "1", "1", "100"])
    assert d2["skewness"] > 1.0, d2
    assert d2["classification"] in ("highly skewed", "heavy-tailed"), d2
    # Insufficient data
    d3 = analyse.distribution(["1", "2"])
    assert d3["skewness"] is None and "insufficient" in d3["classification"]
    # Constant
    d4 = analyse.distribution(["5", "5", "5", "5"])
    assert d4["classification"] == "constant (no spread)"


def test_trend_slope_r2_direction():
    import analyse
    # Rising: y = 2x → slope=2, R²=1
    rising = [("W1", Decimal("2")), ("W2", Decimal("4")), ("W3", Decimal("6")), ("W4", Decimal("8"))]
    t = analyse.trend(rising)
    assert t["n"] == 4 and t["slope"] == 2.0, t
    assert t["r_squared"] == 1.0, t
    assert t["classification"] == "rising", t
    # Falling
    falling = [("W1", Decimal("10")), ("W2", Decimal("8")), ("W3", Decimal("6")), ("W4", Decimal("4"))]
    t2 = analyse.trend(falling)
    assert t2["slope"] == -2.0 and t2["classification"] == "falling", t2
    # Flat (constant)
    flat = [("W1", Decimal("5")), ("W2", Decimal("5")), ("W3", Decimal("5"))]
    t3 = analyse.trend(flat)
    assert t3["classification"] == "flat", t3
    # Insufficient
    t4 = analyse.trend([("W1", Decimal("1")), ("W2", Decimal("2"))])
    assert t4["slope"] is None and "insufficient" in t4["classification"]
    # Noisy / weak (low R²)
    noisy = [("W1", Decimal("1")), ("W2", Decimal("5")), ("W3", Decimal("2")), ("W4", Decimal("4"))]
    t5 = analyse.trend(noisy)
    assert t5["r_squared"] < 0.5, t5


def test_percentile_quantiles():
    import analyse
    vals = ["10", "20", "30", "40", "50", "60", "70", "80", "90", "100"]
    p50 = analyse.percentile(vals, 0.5)
    assert p50["value"] == Decimal("55"), p50
    p90 = analyse.percentile(vals, 0.9)
    assert p90["value"] == Decimal("91"), p90
    p0 = analyse.percentile(vals, 0.0)
    assert p0["value"] == Decimal("10"), p0
    p100 = analyse.percentile(vals, 1.0)
    assert p100["value"] == Decimal("100"), p100
    multi = analyse.percentile(vals, [0.25, 0.75])
    assert multi[0.25] == Decimal("32.5") and multi[0.75] == Decimal("77.5"), multi
    pe = analyse.percentile([], 0.5)
    assert pe["value"] is None and pe["n"] == 0


def test_cohort_retention_matrix():
    import analyse
    header = ["Customer", "Date", "Amount"]
    rows = [
        ["A", "01/01/2026", "100"],
        ["A", "01/02/2026", "50"],
        ["A", "01/03/2026", "30"],
        ["B", "01/01/2026", "200"],
        ["B", "01/03/2026", "40"],
        ["C", "01/02/2026", "150"],
        ["C", "01/03/2026", "60"],
    ]
    ch = analyse.cohort(header, rows, "Customer", "Date", grain="month")
    assert ch["cohorts"] == ["2026-01", "2026-02"], ch["cohorts"]
    assert ch["cohort_sizes"] == [2, 1], ch["cohort_sizes"]
    # Rectangular: all rows padded to max_offset + 1
    assert ch["max_offset"] == 2, ch["max_offset"]
    assert all(len(r) == ch["max_offset"] + 1 for r in ch["matrix"]), ch["matrix"]
    # Jan cohort: offset 0 = 2 entities, offset 1 = 1 (A only), offset 2 = 2 (A+B)
    assert ch["matrix"][0][0] == Decimal(2), ch["matrix"][0]
    assert ch["matrix"][0][1] == Decimal(1), ch["matrix"][0]
    assert ch["matrix"][0][2] == Decimal(2), ch["matrix"][0]
    # Retention (entity-count-based): 2/2=1.0, 1/2=0.5, 2/2=1.0
    assert ch["retention"][0][0] == Decimal(1), ch["retention"][0]
    assert ch["retention"][0][1] == Decimal("0.5"), ch["retention"][0]
    # Feb cohort: 1 entity, padded to 3 columns
    assert ch["matrix"][1][0] == Decimal(1), ch["matrix"][1]
    assert ch["retention"][1][0] == Decimal(1), ch["retention"][1]
    # Value mode: matrix holds value sums, retention still entity-count-based
    chv = analyse.cohort(header, rows, "Customer", "Date", value="Amount", grain="month")
    assert chv["measure"] == "Amount"
    assert chv["matrix"][0][0] == Decimal("300"), chv["matrix"][0]  # A(100)+B(200)
    assert chv["value_matrix"] is not None
    assert chv["value_matrix"][0][0] == Decimal("300"), chv["value_matrix"][0]
    # Retention in value mode is STILL entity-count-based (0–1 fractions)
    assert chv["retention"][0][0] == Decimal(1), chv["retention"][0]  # 2/2, not 300/300
    assert chv["retention"][0][1] == Decimal("0.5"), chv["retention"][0]  # 1/2


def test_correlation_matrix_pairwise():
    import analyse
    header = ["Revenue", "Headcount", "Spend"]
    rows = [["100", "10", "30"], ["150", "12", "40"], ["200", "15", "50"],
            ["250", "18", "60"], ["300", "20", "70"]]
    cm = analyse.correlation_matrix(header, rows, ["Revenue", "Headcount", "Spend"])
    assert cm["n_cols"] == 3
    assert cm["matrix"][0][0] == 1.0
    assert cm["matrix"][0][1] is not None and cm["matrix"][0][1] > 0.95
    assert cm["matrix"][1][0] == cm["matrix"][0][1]  # symmetric
    cm2 = analyse.correlation_matrix(["A", "B"], [["1", "2"], ["3", "4"]], ["A", "B"])
    assert cm2["matrix"][0][1] is None  # only 2 rows
    # M2 fix: junk cell in one column does not shift row pairing
    # Row 2 has junk in col B; without row-wise alignment, col A's [2,4,6,8]
    # would pair with col B's [20,6,8] (shifted). With fix, pairs are (2,20?),
    # (4,junk→skip), (6,6), (8,8) — only rows where BOTH parse.
    cm3 = analyse.correlation_matrix(["A", "B"],
        [["2", "10"], ["4", "junk"], ["6", "6"], ["8", "8"]], ["A", "B"])
    # Only 3 common rows (1,3,4) → correlation computable, not shifted
    assert cm3["matrix"][0][1] is not None


def test_rolling_moving_average():
    import analyse
    series = [("W1", Decimal("10")), ("W2", Decimal("20")), ("W3", Decimal("30")),
              ("W4", Decimal("40")), ("W5", Decimal("50"))]
    r3 = analyse.rolling(series, 3, func="mean")
    assert r3[0] == ("W1", None) and r3[1] == ("W2", None)
    assert r3[2] == ("W3", Decimal("20")), r3[2]
    assert r3[3] == ("W4", Decimal("30")), r3[3]
    assert r3[4] == ("W5", Decimal("40")), r3[4]
    rs = analyse.rolling(series, 2, func="sum")
    assert rs[1] == ("W2", Decimal("30")), rs[1]
    rm = analyse.rolling(series, 3, func="median")
    assert rm[2] == ("W3", Decimal("20")), rm[2]
    r1 = analyse.rolling(series, 1)
    assert r1[0] == ("W1", Decimal("10")), r1[0]


def test_gini_inequality():
    import analyse
    g_equal = analyse.gini(["10", "10", "10", "10"])
    assert g_equal["gini"] == 0.0, g_equal
    assert g_equal["classification"] == "relatively equal", g_equal
    g_unequal = analyse.gini(["0", "0", "0", "1000"])
    assert g_unequal["gini"] > 0.7, g_unequal
    assert g_unequal["classification"] == "extreme inequality", g_unequal
    g_neg = analyse.gini(["100", "-50", "60"])
    assert g_neg["gini"] is None and "negatives" in g_neg["classification"]
    g_one = analyse.gini(["100"])
    assert g_one["gini"] is None and "insufficient" in g_one["classification"]


def test_seasonality_month_and_quarter():
    import analyse
    header = ["Date", "Revenue"]
    rows = [
        ["15/01/2025", "100"], ["15/01/2026", "120"],
        ["15/07/2025", "200"], ["15/07/2026", "220"],
        ["15/10/2025", "300"],
    ]
    sm = analyse.seasonality(header, rows, "Date", value="Revenue", grain="month")
    assert sm["grain"] == "month" and len(sm["seasons"]) == 12
    jan = [s for s in sm["seasons"] if s["season"] == 1][0]
    jul = [s for s in sm["seasons"] if s["season"] == 7][0]
    oct_s = [s for s in sm["seasons"] if s["season"] == 10][0]
    assert jan["count"] == 2 and jan["total"] == Decimal("220"), jan
    assert jan["average"] == Decimal("110"), jan
    assert jul["total"] == Decimal("420"), jul
    assert oct_s["count"] == 1 and oct_s["total"] == Decimal("300"), oct_s
    # Overall average: mean of seasons WITH data (Jan, Jul, Oct = 3 seasons)
    assert sm["overall_average"] == Decimal("940") / Decimal(3), sm["overall_average"]
    assert sm["n_seasons_with_data"] == 3
    # Quarter grain
    sq = analyse.seasonality(header, rows, "Date", value="Revenue", grain="quarter")
    assert sq["grain"] == "quarter" and len(sq["seasons"]) == 4
    q1 = [s for s in sq["seasons"] if s["season"] == 1][0]
    assert q1["total"] == Decimal("220"), q1
    sc = analyse.seasonality(header, rows, "Date", grain="month")
    assert sc["measure"] == "rows"
    assert [s for s in sc["seasons"] if s["season"] == 1][0]["count"] == 2


# --------------------------------------------------------------------------- #
# Standalone runner (no pytest needed)
# --------------------------------------------------------------------------- #
def _run_all():
    tests = sorted((n, f) for n, f in globals().items()
                   if n.startswith("test_") and callable(f))
    failed = 0
    for name, fn in tests:
        try:
            fn()
            print(f"  PASS  {name}")
        except Exception as e:  # noqa: BLE001
            failed += 1
            print(f"  FAIL  {name}: {type(e).__name__}: {e}")
    print(f"\n{len(tests) - failed}/{len(tests)} passed"
          + (f", {failed} FAILED" if failed else ""))
    return 1 if failed else 0


if __name__ == "__main__":
    raise SystemExit(_run_all())
